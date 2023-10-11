from pypdf import PdfReader
from pptx import Presentation

pdf_file = "input.pdf"
pptx_file = "output.pptx"

# 读取PDF文件
reader = PdfReader(pdf_file)
number_of_pages = len(reader.pages)

# 创建幻灯片
presentation = Presentation()

for pageNum in range(number_of_pages):
    pageObj = reader.pages[pageNum]

    text = pageObj.extract_text()
    title = f"Slide {pageNum + 1}"

    # 添加幻灯片并插入文本
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = text

# 保存幻灯片
presentation.save(pptx_file)

print("PDF to PPT conversion completed!")
